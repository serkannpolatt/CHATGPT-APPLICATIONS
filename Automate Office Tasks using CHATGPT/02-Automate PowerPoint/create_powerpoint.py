import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Determine the directory of the script
# Komut dosyasının dizinini belirleyin
script_dir = os.path.dirname(os.path.abspath(__file__))

input_folder = os.path.join(script_dir, 'input')
charts_folder = os.path.join(script_dir, 'charts')
ppt_file = 'financial_data.pptx'

# Create the charts folder if it doesn't exist
# Mevcut değilse, grafikler klasörü oluşturun
if not os.path.exists(charts_folder):
    os.mkdir(charts_folder)

# Create a new PowerPoint presentation
# Yeni bir PowerPoint sunumu oluşturun
prs = Presentation()

# Iterate through all Excel files in the input folder
# Giriş klasöründeki tüm Excel dosyalarını yineleyin
for excel_file in os.listdir(input_folder):
    if not excel_file.endswith('.xlsx'):
        continue

    # Read the financial data from the first worksheet of the Excel file
    # Finansal verileri Excel dosyasının ilk çalışma sayfasından okuyun
    file_path = os.path.join(input_folder, excel_file)
    df = pd.read_excel(file_path, sheet_name=0, usecols="A:P")
    df = df.dropna()

    # Group the data by the "Product" column and sum up the "Sales" column
    # Verileri "Ürün" sütununa göre gruplandırın ve "Satış" sütununu özetleyin
    grouped = df.groupby('Product').sum()['Sales']

    # Create a chart using the seaborn library
    # Seaborn kitaplığını kullanarak bir grafik oluşturun
    sns.barplot(x=grouped.index, y=grouped.values)
    plt.title(excel_file)
    plt.xlabel('Product/Ürün')
    plt.ylabel('Sales/Satış')
    plt.tight_layout()

    # Save the chart to the charts folder
    # Grafiği grafikler klasörüne kaydedin
    chart_file = excel_file.replace('.xlsx', '.png')
    chart_path = os.path.join(charts_folder, chart_file)
    plt.savefig(chart_path)

    # Add a slide to the PowerPoint presentation and insert the chart and title
    # PowerPoint sunumuna bir slayt ekleyin ve grafiği ve başlığı ekleyin
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = excel_file.replace('.xlsx','')

    chart_file = chart_path
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(9)
    height = Inches(6)
    slide.shapes.add_picture(chart_file, left, top, width=width, height=height)

# Save the PowerPoint presentation in the same directory as the script
# PowerPoint sunumunu komut dosyasıyla aynı dizine kaydedin
ppt_path = os.path.join(script_dir, ppt_file)
prs.save(ppt_path)
